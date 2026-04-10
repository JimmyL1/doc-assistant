from pathlib import Path

from langchain_core.documents import Document


def parse_pptx(file_path: str | Path) -> list[Document]:
    from pptx import Presentation

    file_path = Path(file_path)
    prs = Presentation(file_path)

    documents: list[Document] = []

    for slide_index, slide in enumerate(prs.slides):
        text_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    text_parts.append(line)

        if not text_parts:
            continue

        documents.append(Document(
            page_content="\n".join(text_parts),
            metadata={
                "source": file_path.name,
                "file_type": "pptx",
                "page": slide_index,
                "slide": slide_index + 1,
            },
        ))

    if not documents:
        raise ValueError(f"No text could be extracted from '{file_path.name}'")

    return documents
